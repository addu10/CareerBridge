import re
from typing import Dict, List, Any
import google.generativeai as genai
import json
from django.conf import settings
import io
from pdfminer.high_level import extract_text as pdf_extract_text
from docx import Document
import filetype
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Configure Gemini API
genai.configure(
    api_key="AIzaSyC9thMMTST1V85nN8e7EjISEVegtttlBrI"
)
model = genai.GenerativeModel('gemini-1.5-flash')

def analyze_resume(resume_text: str) -> Dict[str, Any]:
    """
    Analyze resume using Gemini AI
    """
    prompt = f"""
    You are an expert resume reviewer with years of experience in talent acquisition and career counseling.
    Your task is to provide a detailed analysis of this resume focusing on its overall quality, content, and effectiveness.

    Resume:
    {resume_text}

    Follow these steps:
    1. Evaluate the overall resume quality and structure
    2. Identify key strengths and accomplishments
    3. Spot areas needing improvement
    4. Provide specific, actionable recommendations
    5. Calculate an overall score based on industry standards

    Return your analysis in this exact JSON format without any additional text or explanation:
    {{
        "score": <number between 0-100>,
        "strengths": [
            <list of specific strengths, achievements, and positive aspects>,
            <include both content and format strengths>,
            <highlight unique selling points>
        ],
        "weaknesses": [
            <list of specific areas needing improvement>,
            <include both content and format issues>,
            <focus on critical issues first>
        ],
        "improvements": [
            <list of specific, actionable recommendations>,
            <prioritize by impact>,
            <include both quick wins and long-term improvements>
        ]
    }}
    """
    
    try:
        response = model.generate_content(prompt)
        if not response.text:
            raise ValueError("Empty response from Gemini API")
            
        # Extract JSON from the response text
        # First, try direct parsing
        try:
            result = json.loads(response.text)
        except json.JSONDecodeError:
            # If direct parsing fails, try to extract JSON from the text
            text = response.text
            # Find JSON-like content between curly braces
            match = re.search(r'\{.*\}', text, re.DOTALL)
            if match:
                try:
                    result = json.loads(match.group(0))
                except:
                    raise ValueError("Could not parse JSON from response")
            else:
                raise ValueError("No JSON content found in response")
        
        # Ensure all required fields are present
        required_fields = ['score', 'strengths', 'weaknesses', 'improvements']
        if not all(field in result for field in required_fields):
            raise ValueError("Missing required fields in AI response")
            
        # Ensure lists are not empty
        if not (result['strengths'] and result['weaknesses'] and result['improvements']):
            raise ValueError("Empty analysis lists in AI response")
            
        return result
        
    except json.JSONDecodeError:
        print("Error: Invalid JSON response from Gemini API")
        return _get_fallback_response("Error processing resume analysis")
    except Exception as e:
        print(f"Error in resume analysis: {str(e)}")
        return _get_fallback_response(str(e))

def _get_fallback_response(error_message: str) -> Dict[str, Any]:
    """Return a structured fallback response when analysis fails"""
    return {
        "score": 0,
        "strengths": ["Could not analyze resume strengths"],
        "weaknesses": ["Error analyzing resume"],
        "improvements": ["Please try submitting your resume again"],
        "error": error_message
    }

def prepare_interview(job_title: str, company: str, resume_text: str) -> Dict[str, Any]:
    """
    Generate interview preparation content based on job details and resume.
    """
    prompt = f"""
    Please prepare interview content for a {job_title} position at {company}.
    Use the candidate's resume to personalize the preparation.
    Return the response in this JSON format:
    {{
        "preparation_content": "General preparation advice",
        "common_questions": [list of common questions],
        "technical_questions": [list of technical questions],
        "behavioral_questions": [list of behavioral questions],
        "tips": [list of interview tips]
    }}

    Resume:
    {resume_text}
    """

    try:
        response = model.generate_content(prompt)
        # Extract JSON from the response text
        # First, try direct parsing
        try:
            result = json.loads(response.text)
        except json.JSONDecodeError:
            # If direct parsing fails, try to extract JSON from the text
            text = response.text
            # Find JSON-like content between curly braces
            match = re.search(r'\{.*\}', text, re.DOTALL)
            if match:
                try:
                    result = json.loads(match.group(0))
                except:
                    raise ValueError("Could not parse JSON from response")
            else:
                raise ValueError("No JSON content found in response")
        
        return result
        
    except Exception as e:
        return {
            'preparation_content': 'Error generating content',
            'common_questions': ['Unable to generate questions'],
            'technical_questions': ['Error occurred'],
            'behavioral_questions': ['Please try again'],
            'tips': ['Service temporarily unavailable']
        }

def generate_career_roadmap(
    skills: List[str],
    career_goals: str,
    preferences: Dict[str, Any] = None
) -> Dict[str, Any]:
    """
    Generate a career roadmap based on current skills and career goals.
    """
    preferences_text = "\nPreferences:\n" + str(preferences) if preferences else ""
    
    prompt = f"""
    Please create a detailed career roadmap based on these details:
    
    Current Skills:
    {', '.join(skills)}
    
    Career Goals:
    {career_goals}
    {preferences_text}
    
    Return the response in this JSON format:
    {{
        "roadmap_content": {{
            "title": "Title for the roadmap",
            "description": "Overall career path description",
            "phases": [
                {{
                    "duration": "Phase duration (e.g., '0-3 months')",
                    "skills": ["Skill 1", "Skill 2", "Skill 3"]
                }}
            ],
            "milestones": [
                {{
                    "title": "Milestone title",
                    "description": "Milestone description"
                }}
            ],
            "skills": [
                {{
                    "category": "Skill category (e.g., 'Technical', 'Soft Skills')",
                    "skills": [
                        {{
                            "name": "Skill name",
                            "priority": "High/Medium/Low",
                            "progress": 0
                        }}
                    ]
                }}
            ],
            "projects": [
                {{
                    "title": "Project title",
                    "description": "Project description",
                    "features": ["Feature 1", "Feature 2"]
                }}
            ],
            "resources": [
                {{
                    "category": "Resource category (e.g., 'Books', 'Courses')",
                    "resources": [
                        {{
                            "title": "Resource title",
                            "description": "Resource description"
                        }}
                    ]
                }}
            ]
        }},
        "milestones": ["Milestone 1", "Milestone 2"],
        "skills_to_acquire": ["Skill 1", "Skill 2"],
        "timeline": {{
            "short_term": ["0-6 month goal 1", "0-6 month goal 2"],
            "medium_term": ["6-18 month goal 1", "6-18 month goal 2"],
            "long_term": ["18+ month goal 1", "18+ month goal 2"]
        }}
    }}
    """

    try:
        print(f"Sending prompt to model: {prompt[:100]}...")
        response = model.generate_content(prompt)
        print(f"Received response from model: {response.text[:100]}...")
        
        # Extract JSON from the response text
        text = response.text
        
        # Try to find JSON content in the response
        json_pattern = r'```json\s*([\s\S]*?)\s*```|```\s*([\s\S]*?)\s*```|\{[\s\S]*\}'
        match = re.search(json_pattern, text, re.DOTALL)
        
        if match:
            # Extract the matched JSON content
            json_str = match.group(1) or match.group(2) or match.group(0)
            
            # Clean up the JSON string
            json_str = json_str.strip()
            if not json_str.startswith('{'):
                # Find the first opening brace
                start_idx = json_str.find('{')
                if start_idx != -1:
                    json_str = json_str[start_idx:]
            
            # Find the last closing brace
            end_idx = json_str.rfind('}')
            if end_idx != -1:
                json_str = json_str[:end_idx+1]
            
            print(f"Extracted JSON string: {json_str[:100]}...")
            
            try:
                result = json.loads(json_str)
                print("Successfully parsed JSON")
                
                # Ensure the roadmap_content is properly structured
                if isinstance(result.get('roadmap_content'), str):
                    # Convert string to structured content if needed
                    print("Converting string roadmap_content to structured object")
                    result['roadmap_content'] = {
                        'title': 'Career Roadmap',
                        'description': result['roadmap_content'],
                        'phases': [],
                        'milestones': [],
                        'skills': [],
                        'projects': [],
                        'resources': []
                    }
                
                print(f"Final result keys: {result.keys()}")
                return result
                
            except json.JSONDecodeError as e:
                print(f"JSON parsing error: {str(e)}")
                raise ValueError(f"Could not parse JSON from response: {str(e)}")
        else:
            print("No JSON-like content found in the response")
            raise ValueError("No JSON content found in response")
        
    except Exception as e:
        return {
            'roadmap_content': 'Error generating roadmap',
            'milestones': ['Unable to generate milestones'],
            'skills_to_acquire': ['Error occurred'],
            'timeline': {
                'short_term': ['Service temporarily unavailable'],
                'medium_term': ['Please try again later'],
                'long_term': ['Error generating timeline']
            }
        }

def generate_roadmap_presentation(roadmap_data: dict, output_path: str) -> str:
    """
    Generate a PowerPoint presentation from roadmap data
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # Create presentation
    prs = Presentation()
    
    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    if isinstance(roadmap_data['roadmap_content'], str):
        title.text = "Career Roadmap"
        subtitle.text = roadmap_data['roadmap_content']
    else:
        title.text = roadmap_data['roadmap_content'].get('title', 'Career Roadmap')
        subtitle.text = roadmap_data['roadmap_content'].get('description', '')

    # Overview slide
    overview_slide = prs.slides.add_slide(prs.slide_layouts[1])
    overview_slide.shapes.title.text = "Overview"
    content = overview_slide.placeholders[1]
    
    overview_text = []
    if isinstance(roadmap_data['roadmap_content'], dict):
        if roadmap_data['roadmap_content'].get('phases'):
            overview_text.append("Phases:")
            for phase in roadmap_data['roadmap_content']['phases']:
                overview_text.append(f"• {phase['duration']}: {', '.join(phase['skills'])}")
    
    content.text = "\n".join(overview_text) if overview_text else "Career development roadmap based on your skills and goals"

    # Skills slide
    skills_slide = prs.slides.add_slide(prs.slide_layouts[1])
    skills_slide.shapes.title.text = "Skills to Acquire"
    skills_content = skills_slide.placeholders[1]
    
    skills_text = []
    if isinstance(roadmap_data['roadmap_content'], dict) and roadmap_data['roadmap_content'].get('skills'):
        for category in roadmap_data['roadmap_content']['skills']:
            skills_text.append(f"\n{category['category']}:")
            for skill in category['skills']:
                priority = skill.get('priority', 'Medium')
                skills_text.append(f"• {skill['name']} (Priority: {priority})")
    elif roadmap_data.get('skills_to_acquire'):
        for skill in roadmap_data['skills_to_acquire']:
            skills_text.append(f"• {skill}")
    
    skills_content.text = "\n".join(skills_text) if skills_text else "No skills defined"

    # Timeline slide
    timeline_slide = prs.slides.add_slide(prs.slide_layouts[1])
    timeline_slide.shapes.title.text = "Development Timeline"
    timeline_content = timeline_slide.placeholders[1]
    
    timeline_text = []
    if roadmap_data.get('timeline'):
        timeline = roadmap_data['timeline']
        if timeline.get('short_term'):
            timeline_text.append("Short-term (0-6 months):")
            for goal in timeline['short_term']:
                timeline_text.append(f"• {goal}")
        if timeline.get('medium_term'):
            timeline_text.append("\nMedium-term (6-18 months):")
            for goal in timeline['medium_term']:
                timeline_text.append(f"• {goal}")
        if timeline.get('long_term'):
            timeline_text.append("\nLong-term (18+ months):")
            for goal in timeline['long_term']:
                timeline_text.append(f"• {goal}")
    
    timeline_content.text = "\n".join(timeline_text) if timeline_text else "No timeline defined"

    # Projects slide
    if isinstance(roadmap_data['roadmap_content'], dict) and roadmap_data['roadmap_content'].get('projects'):
        projects_slide = prs.slides.add_slide(prs.slide_layouts[1])
        projects_slide.shapes.title.text = "Recommended Projects"
        projects_content = projects_slide.placeholders[1]
        
        projects_text = []
        for project in roadmap_data['roadmap_content']['projects']:
            projects_text.append(f"\n{project['title']}:")
            projects_text.append(f"{project['description']}")
            if project.get('features'):
                projects_text.append("Features:")
                for feature in project['features']:
                    projects_text.append(f"• {feature}")
        
        projects_content.text = "\n".join(projects_text)

    # Resources slide
    if isinstance(roadmap_data['roadmap_content'], dict) and roadmap_data['roadmap_content'].get('resources'):
        resources_slide = prs.slides.add_slide(prs.slide_layouts[1])
        resources_slide.shapes.title.text = "Learning Resources"
        resources_content = resources_slide.placeholders[1]
        
        resources_text = []
        for category in roadmap_data['roadmap_content']['resources']:
            resources_text.append(f"\n{category['category']}:")
            for resource in category['resources']:
                resources_text.append(f"• {resource['title']}")
                if resource.get('description'):
                    resources_text.append(f"  {resource['description']}")
        
        resources_content.text = "\n".join(resources_text)

    # Save presentation
    prs.save(output_path)
    return output_path

class ATSReviewer:
    def __init__(self):
        self.model = model
    
    def review_resume(self, resume_text: str) -> Dict[str, Any]:
        """
        Perform an ATS compatibility review of a resume.
        """
        if self.model is None:
            return {
                'score': 0,
                'feedback': ['AI model is not available. Please check your API configuration.']
            }
            
        prompt = f"""
        Please review this resume for ATS (Applicant Tracking System) compatibility.
        Provide a detailed analysis in this exact JSON format without any additional text:
        {{
            "score": (0-100 score based on ATS compatibility),
            "feedback": [list of specific feedback points]
        }}

        Resume:
        {resume_text}
        """

        try:
            response = self.model.generate_content(prompt)
            if not response.text:
                return {
                    'score': 0,
                    'feedback': ['Empty response from AI model']
                }
                
            # Extract JSON from the response text
            # First, try direct parsing
            try:
                result = json.loads(response.text)
            except json.JSONDecodeError:
                # If direct parsing fails, try to extract JSON from the text
                text = response.text
                # Find JSON-like content between curly braces
                match = re.search(r'\{.*\}', text, re.DOTALL)
                if match:
                    try:
                        result = json.loads(match.group(0))
                    except:
                        print("Could not parse JSON from response")
                        return {
                            'score': 0,
                            'feedback': ['Error processing ATS compatibility analysis']
                        }
                else:
                    print("No JSON content found in response")
                    return {
                        'score': 0,
                        'feedback': ['Error processing ATS compatibility analysis']
                    }
            
            # Validate required fields
            if not all(field in result for field in ['score', 'feedback']):
                print("Missing required fields in AI response")
                return {
                    'score': 0,
                    'feedback': ['Incomplete analysis from AI model']
                }
                
            return result
            
        except Exception as e:
            print(f"Error in ATS analysis: {str(e)}")
            return {
                'score': 0,
                'feedback': ['Error analyzing ATS compatibility']
            }

    def extract_keywords(self, text: str) -> List[str]:
        """
        Extract important keywords from text using AI.
        """
        prompt = f"""
        Please extract important keywords from this text that would be relevant
        for ATS systems. Return only a Python list of lowercase strings.
        
        Text:
        {text}
        """

        try:
            response = self.model.generate_content(prompt)
            # Extract JSON from the response text
            # First, try direct parsing
            try:
                result = json.loads(response.text)
            except json.JSONDecodeError:
                # If direct parsing fails, try to extract JSON from the text
                text = response.text
                # Find JSON-like content between curly braces
                match = re.search(r'\{.*\}', text, re.DOTALL)
                if match:
                    try:
                        result = json.loads(match.group(0))
                    except:
                        raise ValueError("Could not parse JSON from response")
                else:
                    raise ValueError("No JSON content found in response")
            
            keywords = result
            return [k.lower() for k in keywords if isinstance(k, str)]
            
        except Exception as e:
            return []

def extract_text_from_file(file):
    """
    Extract text from PDF, DOC, DOCX, or RTF files
    """
    # Read the file content
    file_content = file.read()
    file.seek(0)  # Reset file pointer
    
    # Detect file type using filetype
    kind = filetype.guess(file_content)
    if kind is None:
        # If filetype can't detect it, try checking the file extension
        file_name = str(file.name).lower()
        if file_name.endswith('.rtf'):
            file_type = 'rtf'
        else:
            raise ValueError("Could not determine file type")
    else:
        file_type = kind.mime
    
    if 'pdf' in file_type:
        # PDF file
        print(pdf_extract_text(io.BytesIO(file_content)))
        return pdf_extract_text(io.BytesIO(file_content))
    elif 'msword' in file_type or 'vnd.openxmlformats-officedocument.wordprocessingml.document' in file_type or str(file.name).lower().endswith('.docx'):
        # DOC/DOCX file
        doc = Document(io.BytesIO(file_content))
        return '\n'.join([paragraph.text for paragraph in doc.paragraphs])
    elif file_type == 'rtf' or str(file.name).lower().endswith('.rtf'):
        # RTF file - convert to plain text by stripping RTF markup
        text = file_content.decode('utf-8', errors='ignore')
        # Basic RTF cleanup (this is a simple approach, might need improvement)
        text = re.sub(r'[\\\{\}]|\\\w+\s?', '', text)
        return text
    else:
        raise ValueError(f"Unsupported file type. Please upload a PDF, DOC, DOCX, or RTF file.")